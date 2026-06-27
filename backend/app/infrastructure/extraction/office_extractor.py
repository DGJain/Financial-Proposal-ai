"""DOCX / PPTX extraction via python-docx and python-pptx.

Lazy-imports the Office libraries inside ``extract``. Word documents have no fixed
page model, so paragraphs/tables are emitted as a single logical page; PowerPoint
maps one slide to one page. Both converge on the shared ``ExtractedDocument``.
"""

from __future__ import annotations

from io import BytesIO

from app.domain.documents.enums import FileType
from app.domain.ingestion.extracted import (
    ExtractedDocument,
    ExtractedPage,
    ExtractedTable,
    TextBlock,
)


class OfficeExtractor:
    """``ExtractorPort`` adapter for DOCX and PPTX sources."""

    def supports(self, file_type: FileType) -> bool:
        return file_type in (FileType.DOCX, FileType.PPTX)

    async def extract(self, data: bytes, *, file_type: FileType) -> ExtractedDocument:
        if file_type is FileType.DOCX:
            return self._extract_docx(data)
        return self._extract_pptx(data)

    def _extract_docx(self, data: bytes) -> ExtractedDocument:
        import docx  # lazy

        document = docx.Document(BytesIO(data))
        blocks = tuple(
            TextBlock(text=p.text.strip()) for p in document.paragraphs if p.text.strip()
        )
        tables: list[ExtractedTable] = []
        for index, table in enumerate(document.tables):
            rows = tuple(
                tuple(cell.text.strip() for cell in row.cells) for row in table.rows
            )
            if rows:
                tables.append(
                    ExtractedTable(
                        table_id=f"docx-t{index}",
                        rows=rows,
                        has_total_row=_looks_like_total(rows[-1]),
                    )
                )
        page = ExtractedPage(page_number=1, text_blocks=blocks, tables=tuple(tables))
        return ExtractedDocument(file_type=FileType.DOCX, pages=(page,))

    def _extract_pptx(self, data: bytes) -> ExtractedDocument:
        from pptx import Presentation  # lazy

        presentation = Presentation(BytesIO(data))
        pages: list[ExtractedPage] = []
        for index, slide in enumerate(presentation.slides, start=1):
            blocks: list[TextBlock] = []
            tables: list[ExtractedTable] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    blocks.append(TextBlock(text=shape.text_frame.text.strip()))
                if getattr(shape, "has_table", False):
                    rows = tuple(
                        tuple(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    )
                    if rows:
                        tables.append(
                            ExtractedTable(
                                table_id=f"slide{index}-t{len(tables)}",
                                rows=rows,
                                has_total_row=_looks_like_total(rows[-1]),
                            )
                        )
            pages.append(
                ExtractedPage(page_number=index, text_blocks=tuple(blocks), tables=tuple(tables))
            )
        return ExtractedDocument(file_type=FileType.PPTX, pages=tuple(pages))


def _looks_like_total(row: tuple[str, ...]) -> bool:
    joined = " ".join(row).lower()
    return any(marker in joined for marker in ("total", "net", "sum"))
